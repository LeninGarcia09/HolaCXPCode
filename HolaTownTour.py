import requests
import wikipedia
from pptx import Presentation
from pptx.util import Inches
import os
# Add tkinter for GUI
import tkinter as tk
from tkinter import messagebox
from tkinter import simpledialog
# Add OpenAI for AI enrichment
import openai


# Unsplash API setup
UNSPLASH_ACCESS_KEY = "YOUR_UNSPLASH_ACCESS_KEY"  # <-- replace with your key

# OpenAI API setup
OPENAI_API_KEY = "YOUR_OPENAI_API_KEY"  # <-- replace with your key
openai.api_key = OPENAI_API_KEY
# Fun facts from Wikidata
def get_fun_facts(hometown, max_facts=3):
    # ...existing code...

# AI enrichment function
def get_ai_enriched_content(hometown, traditions, remarks):
    prompt = (
        f"Create a fun and engaging summary for the town {hometown}, "
        f"including its history, famous places, local foods, and a few fun facts. "
        f"Traditions: {', '.join(traditions)}. "
        f"Remarks: {', '.join(remarks)}. "
        f"Make it interesting and suitable for a presentation slide."
    )
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=500
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return "AI enrichment is currently unavailable. Please check your API key or network connection."


def get_wikipedia_summary(hometown):
    # Try Wikipedia first
    try:
        return wikipedia.summary(hometown, sentences=3)
    except Exception:
        pass
    # Try Wikidata
    try:
        url = f"https://www.wikidata.org/w/api.php"
        params = {
            "action": "wbsearchentities",
            "search": hometown,
            "language": "en",
            "format": "json"
        }
        response = requests.get(url, params=params)
        data = response.json()
        if data.get("search"):
            description = data["search"][0].get("description", None)
            if description:
                return f"{hometown}: {description}"
    except Exception:
        pass
    # Fallback
    return f"{hometown} is known for its rich culture and history."

# Fun facts from Wikidata
def get_fun_facts(hometown, max_facts=3):
    facts = []
    try:
        # Search for the town entity
        url = "https://www.wikidata.org/w/api.php"
        params = {
            "action": "wbsearchentities",
            "search": hometown,
            "language": "en",
            "format": "json"
        }
        response = requests.get(url, params=params)
        data = response.json()
        if data.get("search"):
            entity_id = data["search"][0]["id"]
            # Get claims (properties) for the entity
            entity_url = f"https://www.wikidata.org/wiki/Special:EntityData/{entity_id}.json"
            entity_resp = requests.get(entity_url)
            entity_data = entity_resp.json()
            claims = entity_data["entities"][entity_id]["claims"]
            # Example: population, inception, country, etc.
            if "P1082" in claims:  # population
                pop = claims["P1082"][0]["mainsnak"]["datavalue"]["value"]
                facts.append(f"Population: {pop}")
            if "P571" in claims:  # inception
                inception = claims["P571"][0]["mainsnak"]["datavalue"]["value"]["time"]
                facts.append(f"Founded: {inception}")
            if "P17" in claims:  # country
                country_id = claims["P17"][0]["mainsnak"]["datavalue"]["value"]["id"]
                facts.append(f"Country: https://www.wikidata.org/wiki/{country_id}")
    except Exception:
        pass
    if not facts:
        facts.append("Did you know? This town has many interesting stories!")
    return facts[:max_facts]


def get_unsplash_image(hometown, count=1, save_dir="images"):
    url = f"https://api.unsplash.com/search/photos"
    params = {"query": hometown, "client_id": UNSPLASH_ACCESS_KEY, "per_page": count}
    response = requests.get(url, params=params)
    data = response.json()

    if not os.path.exists(save_dir):
        os.makedirs(save_dir)

    image_files = []
    for i, result in enumerate(data.get("results", [])):
        img_url = result["urls"]["regular"]
        img_data = requests.get(img_url).content
        img_path = os.path.join(save_dir, f"{hometown}_{i}.jpg")
        with open(img_path, "wb") as f:
            f.write(img_data)
        image_files.append(img_path)

    return image_files




def create_hometown_ppt(hometown, traditions, remarks, filename="hometown_presentation.pptx"):
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"Welcome to {hometown}"
    slide.placeholders[1].text = "A journey through its traditions and culture"

    # About slide (Wikipedia/Wikidata summary)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"About {hometown}"
    slide.placeholders[1].text = get_wikipedia_summary(hometown)

    # AI Enriched Content slide
    ai_content = get_ai_enriched_content(hometown, traditions, remarks)
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"Discover {hometown}"
    slide.placeholders[1].text = ai_content

    # Fun Facts slide
    fun_facts = get_fun_facts(hometown)
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"Fun Facts about {hometown}"
    slide.placeholders[1].text = "\n".join([f"• {fact}" for fact in fun_facts])

    # Traditions slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"Traditions of {hometown}"
    slide.placeholders[1].text = "\n".join([f"• {t}" for t in traditions])

    # Remarks slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Highlights & Remarks"
    slide.placeholders[1].text = "\n".join([f"• {r}" for r in remarks])

    # Image slide (from Unsplash)
    images = get_unsplash_image(hometown, count=2)
    for img_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        slide.shapes.title.text = f"{hometown} in Pictures"
        left = Inches(1)
        top = Inches(2)
        height = Inches(4)
        slide.shapes.add_picture(img_path, left, top, height=height)

    # Closing slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"Thank You for Visiting {hometown}"
    slide.placeholders[1].text = "We hope you enjoyed learning about this hometown!"

    prs.save(filename)
    print(f"Presentation saved as {filename}")



# GUI for user input
def run_gui():
    def submit():
        hometown = entry_town.get().strip()
        traditions = text_traditions.get("1.0", tk.END).strip().splitlines()
        remarks = text_remarks.get("1.0", tk.END).strip().splitlines()
        if not hometown:
            messagebox.showerror("Error", "Please enter the name of the town.")
            return
        if not traditions or all(t.strip() == "" for t in traditions):
            messagebox.showerror("Error", "Please enter at least one tradition.")
            return
        if not remarks or all(r.strip() == "" for r in remarks):
            messagebox.showerror("Error", "Please enter at least one remark.")
            return
        create_hometown_ppt(hometown, traditions, remarks)
        messagebox.showinfo("Success", f"Presentation for {hometown} created!")

    root = tk.Tk()
    root.title("Hometown Tour Presentation Generator")

    tk.Label(root, text="Town Name:").grid(row=0, column=0, sticky="w", padx=5, pady=5)
    entry_town = tk.Entry(root, width=40)
    entry_town.grid(row=0, column=1, padx=5, pady=5)

    tk.Label(root, text="Traditions (one per line):").grid(row=1, column=0, sticky="nw", padx=5, pady=5)
    text_traditions = tk.Text(root, width=40, height=5)
    text_traditions.grid(row=1, column=1, padx=5, pady=5)

    tk.Label(root, text="Important Remarks (one per line):").grid(row=2, column=0, sticky="nw", padx=5, pady=5)
    text_remarks = tk.Text(root, width=40, height=5)
    text_remarks.grid(row=2, column=1, padx=5, pady=5)

    submit_btn = tk.Button(root, text="Create Presentation", command=submit)
    submit_btn.grid(row=3, column=0, columnspan=2, pady=10)

    root.mainloop()


if __name__ == "__main__":
    run_gui()